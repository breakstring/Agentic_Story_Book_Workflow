from typing import Annotated
from pptx import Presentation
from pptx.util import Inches
import os
from PIL import Image
from .utils import get_storyboard_by_story_id

def create_pptx(story_id: Annotated[str,'story id']):

    output_directory= "./output/"+story_id
    storyboard = get_storyboard_by_story_id(story_id)

    prs = Presentation()
    
    # 获取所有子目录
    subdirs = sorted([d for d in os.listdir(output_directory) if os.path.isdir(os.path.join(output_directory, d))], key=int)
    
    for subdir in subdirs:
        subdir_path = os.path.join(output_directory, subdir)
        image_path = os.path.join(subdir_path, 'image.png')
        audio_path = os.path.join(subdir_path, 'voice.mp3')
        
        # 添加新的幻灯片
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 获取当前索引的 storyboard 内容
        storyboard_item = next((item for item in storyboard if item["Index"] == subdir), None)
        if storyboard_item:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = storyboard_item["StoryContent"]
        
        # 打开图片并获取尺寸
        img = Image.open(image_path)
        img_width, img_height = img.size
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # 计算缩放比例
        scale = min((slide_width - Inches(1)) / img_width, (slide_height - Inches(1)) / img_height)
        
        # 计算新的图片尺寸
        new_width = int(img_width * scale)
        new_height = int(img_height * scale)
        
        # 计算图片位置（居中）
        left = int((slide_width - new_width) / 2)
        top = int((slide_height - new_height) / 2)
        
        # 添加图片
        pic = slide.shapes.add_picture(image_path, left, top, new_width, new_height)
        
        # 添加音频
        audio_left = Inches(0)
        audio_top = Inches(0)
        audio_width = Inches(1)
        audio_height = Inches(1)
        slide.shapes.add_movie(audio_path, audio_left, audio_top, audio_width, audio_height)
    
    # 保存 PPTX 文件
    prs.save(os.path.join(output_directory, 'output.pptx'))
    